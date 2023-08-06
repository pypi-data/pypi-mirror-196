from typing import cast

import daiquiri
from pptx.chart.axis import DateAxis
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slide

from dbnomics_pptx_tools.charts.series import update_series
from dbnomics_pptx_tools.data_loader import ShapeDataLoader
from dbnomics_pptx_tools.metadata import ChartSpec
from dbnomics_pptx_tools.pptx_copy import copy_shape_properties
from dbnomics_pptx_tools.repo import SeriesRepo
from dbnomics_pptx_tools.xml_utils import remove_element

from .chart_data import build_category_chart_data
from .data_labels import update_data_labels

logger = daiquiri.getLogger(__name__)


def recreate_chart(chart_shape: GraphicFrame, *, chart_data: CategoryChartData, slide: Slide) -> Chart:
    chart = cast(Chart, chart_shape.chart)
    remove_element(chart_shape.element)
    new_chart_shape = cast(
        GraphicFrame,
        cast(SlideShapes, slide.shapes).add_chart(
            chart.chart_type, chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height, chart_data
        ),
    )
    copy_shape_properties(chart_shape, new_chart_shape)
    logger.debug("The chart was recreated")
    return cast(Chart, new_chart_shape.chart)


def update_chart(chart_shape: GraphicFrame, *, chart_spec: ChartSpec, repo: SeriesRepo, slide: Slide) -> None:
    chart = cast(Chart, chart_shape.chart)
    if not isinstance(chart.category_axis, DateAxis):
        raise NotImplementedError()

    chart_series_df = ShapeDataLoader(repo).load_shape_df(chart_spec)
    pivoted_df = chart_series_df.pivot(index="period", columns="series_id", values="value")

    chart_data = build_category_chart_data(chart_spec, pivoted_df=pivoted_df)

    try:
        chart.replace_data(chart_data)
    except ValueError:
        chart = recreate_chart(chart_shape, chart_data=chart_data, slide=slide)
    else:
        logger.debug("Chart data was replaced")

    update_data_labels(chart, chart_spec=chart_spec, pivoted_df=pivoted_df)
    update_series(chart, chart_spec=chart_spec)
